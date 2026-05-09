from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette (Refined)
BG_DARK = RGBColor(0x0F, 0x0F, 0x1A)
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x6C, 0x63, 0xFF)      # Electric Purple
ACCENT2 = RGBColor(0x00, 0xD2, 0xFF)     # Cyan
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB0, 0xC0)
SUCCESS = RGBColor(0x00, 0xE6, 0x76)     # Green
WARN = RGBColor(0xFF, 0x9F, 0x43)        # Orange
PINK = RGBColor(0xFF, 0x6B, 0x9D)        # Pink/Red
IVORY = RGBColor(0xDF, 0xDF, 0xEF)

def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def set_slide_transition(slide, type="fade"):
    """
    Sets slide transition by manipulating the slide XML directly.
    """
    slide_element = slide.element
    # Find or create transitions element
    if type == "fade":
        trans_xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:fade/></p:transition>'
    elif type == "push":
        # Push from right is 'l' (leftward push)
        trans_xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:push dir="l"/></p:transition>'
    elif type == "wipe":
        trans_xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:wipe/></p:transition>'
    else:
        trans_xml = '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:fade/></p:transition>'
    
    new_element = parse_xml(trans_xml)
    slide_element.insert(len(slide_element), new_element)

def add_shape(slide, left, top, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_accent_bar(slide, left, top, w, h, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_bullet_text(slide, left, top, w, h, items, size=16, color=LIGHT_GRAY, icon="▸"):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{icon} {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
    return txBox

# ─── SLIDE 1: TITLE ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
set_slide_transition(slide, "fade")
add_circle(slide, Inches(-1), Inches(-1), Inches(4), ACCENT)
add_circle(slide, Inches(10.5), Inches(4.5), Inches(5), ACCENT2)
add_shape(slide, Inches(1.5), Inches(1.2), Inches(10.3), Inches(5), BG_CARD)
add_accent_bar(slide, Inches(5.5), Inches(1.8), Inches(2.3), Pt(5), ACCENT)
add_text(slide, Inches(2), Inches(2.1), Inches(9.3), Inches(1.2), "AuraCRM", 56, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(3.2), Inches(9.3), Inches(0.8), "Voice-Driven AI Sales Intelligence", 30, ACCENT2, False, PP_ALIGN.CENTER)
add_accent_bar(slide, Inches(5.8), Inches(4.1), Inches(1.7), Pt(3), ACCENT2)
add_text(slide, Inches(2), Inches(4.4), Inches(9.3), Inches(1), "Transforming CRM Pipelines with Intelligent, Real-Time Conversations", 18, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# ─── SLIDE 2: THE VOICE-FIRST REVOLUTION ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
set_slide_transition(slide, "push")
add_accent_bar(slide, Inches(0.8), Inches(0.6), Inches(0.15), Inches(0.6), ACCENT)
add_text(slide, Inches(1.2), Inches(0.55), Inches(10), Inches(0.7), "The Voice-First Revolution", 36, WHITE, True)
add_text(slide, Inches(1.2), Inches(1.3), Inches(11), Inches(0.5), "\"The Whole Website Works With Voice\" — A Hands-Free CRM Experience", 20, ACCENT2)

add_shape(slide, Inches(0.8), Inches(2.1), Inches(11.7), Inches(4.8), BG_CARD)
items = [
    "Omni-Directional Voice Control: Navigate menus, open files, and search contacts with natural speech.",
    "Active Listening Mode: Aura parses real-time intent from your conversations.",
    "Hands-Free Pipeline Management: Move deals and update statuses without clicking.",
    "Integrated Voice Assistant: Toggles on/off with visual 'Pulse' animations.",
    "Real-Time Speech-to-Intent: Uses advanced NLP to translate voice into platform actions.",
    "Accessibility & Speed: Empowering sales reps to work 2X faster than manual input."
]
add_bullet_text(slide, Inches(1.3), Inches(2.6), Inches(10.7), Inches(4), items, 18, LIGHT_GRAY, "🎙️")

# ─── SLIDE 3: AURA ASSISTANT (AI CHATBOT) ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
set_slide_transition(slide, "push")
add_accent_bar(slide, Inches(0.8), Inches(0.6), Inches(0.15), Inches(0.6), ACCENT)
add_text(slide, Inches(1.2), Inches(0.55), Inches(10), Inches(0.7), "Aura Assistant: Your AI Partner", 36, WHITE, True)

add_shape(slide, Inches(0.8), Inches(1.8), Inches(6), Inches(5), BG_CARD)
add_text(slide, Inches(1.3), Inches(2.1), Inches(5), Inches(0.5), "AI Chat & Voice Synergy", 22, ACCENT, True)
items = [
    "Context-Aware Conversations: Understands which customer you are viewing.",
    "Dual-Mode Interaction: Switch seamlessly between typing and talking.",
    "Text-to-Speech (TTS): Aura replies with natural, premium female voices (Samantha/Victoria).",
    "Product Interest Analysis: Deep dive into customer needs via AI dialogue.",
    "Interactive UI: Glassmorphism window with 'Typing...' indicators."
]
add_bullet_text(slide, Inches(1.3), Inches(2.8), Inches(5), Inches(3.5), items, 16, LIGHT_GRAY, "💬")

add_shape(slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(5), BG_CARD)
add_text(slide, Inches(7.7), Inches(2.1), Inches(4.3), Inches(0.5), "Natural Voice Feedback", 22, PINK, True)
add_text(slide, Inches(7.7), Inches(2.8), Inches(4.3), Inches(1.5), "\"I see you're looking at TechFlow Inc. Would you like me to draft an follow-up or schedule a demo?\"", 18, WHITE, False, PP_ALIGN.LEFT, "Calibri")
add_accent_bar(slide, Inches(7.2), Inches(1.8), Inches(5.3), Pt(4), PINK)

# ─── SLIDE 4: AUTO EMAIL GENERATION ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
set_slide_transition(slide, "push")
add_accent_bar(slide, Inches(0.8), Inches(0.6), Inches(0.15), Inches(0.6), ACCENT)
add_text(slide, Inches(1.2), Inches(0.55), Inches(10), Inches(0.7), "Smart Email Generator", 36, WHITE, True)

add_shape(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5), BG_CARD)
add_accent_bar(slide, Inches(0.8), Inches(1.8), Inches(11.7), Pt(4), SUCCESS)
add_text(slide, Inches(1.3), Inches(2.1), Inches(10), Inches(0.5), "Context-Driven Automation", 22, SUCCESS, True)

items = [
    "Highly Personalized Content: AI drafts emails tailored to customer interest and history.",
    "Tone Selection: Choose between Professional, Friendly, Urgent, or Consultative styles.",
    "Automated Generation: Fresh content generated in 2 seconds via AI Integration.",
    "One-Click Dispatch: Integrated 'Send' button with Nodemailer server backend.",
    "Smart Scheduling: Choose a future date/time to send emails automatically.",
    "Contextual Goal Setting: High-intent outreach, follow-ups, or re-engaging ghosted leads."
]
add_bullet_text(slide, Inches(1.3), Inches(2.8), Inches(10.7), Inches(4), items, 18, LIGHT_GRAY, "📧")

# ─── SLIDE 5: DEAL SUCCESS & SALES INSIGHTS ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
set_slide_transition(slide, "push")
add_accent_bar(slide, Inches(0.8), Inches(0.6), Inches(0.15), Inches(0.6), ACCENT)
add_text(slide, Inches(1.2), Inches(0.55), Inches(10), Inches(0.7), "Sales Insights Engine", 36, WHITE, True)

# Metrics bar
add_shape(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.2), BG_CARD)
add_text(slide, Inches(1.2), Inches(2.0), Inches(3.5), Inches(0.8), "Pipeline Health: 92/100", 22, SUCCESS, True, PP_ALIGN.CENTER)
add_text(slide, Inches(4.9), Inches(2.0), Inches(3.5), Inches(0.8), "Predicted Win Rate: 68%", 22, ACCENT2, True, PP_ALIGN.CENTER)
add_text(slide, Inches(8.6), Inches(2.0), Inches(3.5), Inches(0.8), "AI Recommendations: High", 22, WARN, True, PP_ALIGN.CENTER)

add_shape(slide, Inches(0.8), Inches(3.3), Inches(5.7), Inches(3.5), BG_CARD)
add_text(slide, Inches(1.3), Inches(3.5), Inches(4.7), Inches(0.5), "AI Priority Actions", 22, WARN, True)
items = ["Detects hesitant decision makers via AI Emotion Analysis.", "Identifies technical review completion for proposal generation.", "Flags high-value deals requiring immediate outreach."]
add_bullet_text(slide, Inches(1.3), Inches(4.1), Inches(4.7), Inches(2.5), items, 15, LIGHT_GRAY, "⚡")

add_shape(slide, Inches(6.8), Inches(3.3), Inches(5.7), Inches(3.5), BG_CARD)
add_text(slide, Inches(7.3), Inches(3.5), Inches(4.7), Inches(0.5), "Global Risk Alerts", 22, PINK, True)
items = ["Competitor Surge: Real-time detection of competitor mentions.", "Engagement Drop: Automatic tracking of response time delays.", "AI Sales Coach: Actionable tips to combat risks."]
add_bullet_text(slide, Inches(7.3), Inches(4.1), Inches(4.7), Inches(2.5), items, 15, LIGHT_GRAY, "🚨")

# ─── SLIDE 6: PREMIUM UI / UX AESTHETICS ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
set_slide_transition(slide, "fade")
add_accent_bar(slide, Inches(0.8), Inches(0.6), Inches(0.15), Inches(0.6), ACCENT)
add_text(slide, Inches(1.2), Inches(0.55), Inches(10), Inches(0.7), "Clean & Premium UI Design", 36, WHITE, True)
add_text(slide, Inches(1.2), Inches(1.3), Inches(11), Inches(0.5), "A Visual Experience That Feels 'Alive'", 20, ACCENT2)

cards = [
    ("💎", "Glassmorphism", "Deep shadows, blurred surfaces, and layered depth.", ACCENT),
    ("✨", "Micro-Animations", "Smooth fades, pulses, and transitions.", ACCENT2),
    ("🎨", "Harmonious Palette", "Custom Charcoal/Ivory/Accent theme.", PINK),
    ("📱", "Responsive Layout", "Spacious and adaptable interface.", SUCCESS)
]
for i, (icon, title, desc, col) in enumerate(cards):
    x = Inches(0.8 + (i % 2) * 6.1)
    y = Inches(2.1 + (i // 2) * 2.6)
    add_shape(slide, x, y, Inches(5.7), Inches(2.2), BG_CARD)
    add_accent_bar(slide, x, y, Inches(0.15), Inches(2.2), col)
    add_text(slide, x + Inches(0.5), y + Inches(0.3), Inches(4.8), Inches(0.5), f"{icon}  {title}", 20, WHITE, True)
    add_text(slide, x + Inches(0.5), y + Inches(1.0), Inches(4.8), Inches(1), desc, 16, LIGHT_GRAY)

# ─── SLIDE 7: SMART SCHEDULER & MEET ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
set_slide_transition(slide, "push")
add_accent_bar(slide, Inches(0.8), Inches(0.6), Inches(0.15), Inches(0.6), ACCENT)
add_text(slide, Inches(1.2), Inches(0.55), Inches(10), Inches(0.7), "Smart Follow-Up Scheduler", 36, WHITE, True)

add_shape(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5), BG_CARD)
add_accent_bar(slide, Inches(0.8), Inches(1.8), Inches(11.7), Pt(4), ACCENT2)
items = [
    "Custom Date & Time Selection: Intuitive meeting picker.",
    "Auto Google Meet Links: Unique 'meet.google.com' links generated instantly.",
    "Nodemailer Hub: Sends official invites directly to client inboxes.",
    "Calendar Display: Visual confirmation of scheduled activities.",
    "Seamless Backend: Live integration between Frontend, Server, and Email."
]
add_bullet_text(slide, Inches(1.3), Inches(2.6), Inches(10.7), Inches(4), items, 18, LIGHT_GRAY, "📅")

# ─── SLIDE 8: CUSTOMER LOBBY (CRM) ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
set_slide_transition(slide, "push")
add_accent_bar(slide, Inches(0.8), Inches(0.6), Inches(0.15), Inches(0.6), ACCENT)
add_text(slide, Inches(1.2), Inches(0.55), Inches(10), Inches(0.7), "Customer Lobby & CRM Database", 36, WHITE, True)

add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5), BG_CARD)
add_text(slide, Inches(1.3), Inches(2.1), Inches(4.7), Inches(0.5), "Live Status Tracking", 22, ACCENT, True)
items = ["Real-time badges: Active, Pending, In Pipeline.", "Global search and filtering by interest.", "Quick-action icons for Email/Proposal/Meeting."]
add_bullet_text(slide, Inches(1.3), Inches(2.8), Inches(4.7), Inches(3.5), items, 16, LIGHT_GRAY, "👥")

add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5), BG_CARD)
add_text(slide, Inches(7.3), Inches(2.1), Inches(4.7), Inches(0.5), "Database Management", 22, SUCCESS, True)
items = ["MongoDB Integration: Persistent data storage.", "Direct Deletion: Secure API-based customer removal.", "Safety Confirmation: Guardrails for critical data actions."]
add_bullet_text(slide, Inches(7.3), Inches(2.8), Inches(4.7), Inches(3.5), items, 16, LIGHT_GRAY, "💾")

# ─── SLIDE 9: GLOBAL INBOX WITH AI ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
set_slide_transition(slide, "push")
add_accent_bar(slide, Inches(0.8), Inches(0.6), Inches(0.15), Inches(0.6), ACCENT)
add_text(slide, Inches(1.2), Inches(0.55), Inches(10), Inches(0.7), "Global Inbox & Communications", 36, WHITE, True)

add_shape(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5), BG_CARD)
add_accent_bar(slide, Inches(0.8), Inches(1.8), Inches(11.7), Pt(4), WARN)
items = [
    "Contextual Urgency Badges: Auto-flags High/Medium/Low priority emails.",
    "Integrated Date Filtering: Built-in calendar for communication history.",
    "Inbox/Sent Toggle: Switch between client messages and automated invites.",
    "Communication Audit: Full transparency into all client touchpoints.",
    "Real-time Updates: Fetching latest correspondence from the server."
]
add_bullet_text(slide, Inches(1.3), Inches(2.6), Inches(10.7), Inches(4), items, 18, LIGHT_GRAY, "📥")

# ─── SLIDE 10: PROPOSAL GENERATOR ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
set_slide_transition(slide, "push")
add_accent_bar(slide, Inches(0.8), Inches(0.6), Inches(0.15), Inches(0.6), ACCENT)
add_text(slide, Inches(1.2), Inches(0.55), Inches(10), Inches(0.7), "Account Summarizer & Proposals", 36, WHITE, True)

add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5), BG_CARD)
add_text(slide, Inches(1.3), Inches(2.1), Inches(4.7), Inches(0.5), "Account Summarizer", 22, ACCENT2, True)
items = ["Generates concise briefs of client needs.", "Highlights budget and timeline for reps.", "Allows for 1-click meeting scheduling."]
add_bullet_text(slide, Inches(1.3), Inches(2.8), Inches(4.7), Inches(3.5), items, 16, LIGHT_GRAY, "📄")

add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5), BG_CARD)
add_text(slide, Inches(7.3), Inches(2.1), Inches(4.7), Inches(0.5), "Tailored Proposals", 22, PINK, True)
items = ["Itemized software requirement analysis.", "Automated document generation.", "Customized professional PDF-ready exports."]
add_bullet_text(slide, Inches(7.3), Inches(2.8), Inches(4.7), Inches(3.5), items, 16, LIGHT_GRAY, "📊")

# ─── SLIDE 11: INTERACTIVE DASHBOARD HUB ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
set_slide_transition(slide, "fade")
add_accent_bar(slide, Inches(0.8), Inches(0.6), Inches(0.15), Inches(0.6), ACCENT)
add_text(slide, Inches(1.2), Inches(0.55), Inches(10), Inches(0.7), "Interactive Dashboard Hub", 36, WHITE, True)

add_shape(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5), BG_CARD)
add_accent_bar(slide, Inches(0.8), Inches(1.8), Inches(11.7), Pt(4), ACCENT)
items = [
    "Top Deals Tracker: Visualizes high-value opportunities.",
    "Dynamic Charting: Recharts-powered pipeline health viz.",
    "Real-time Alert Feed: 'Competitor Surge', 'Engagement Drop' notifications.",
    "One-Click Actions: Direct links to Chat, Emails, and Insights.",
    "Global Performance Metrics: Revenue and deal velocity tracking."
]
add_bullet_text(slide, Inches(1.3), Inches(2.6), Inches(10.7), Inches(4), items, 18, LIGHT_GRAY, "🖥️")

# ─── SLIDE 12: TECH STACK ECOSYSTEM ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
set_slide_transition(slide, "push")
add_accent_bar(slide, Inches(0.8), Inches(0.6), Inches(0.15), Inches(0.6), ACCENT)
add_text(slide, Inches(1.2), Inches(0.55), Inches(10), Inches(0.7), "Technology Stack Ecosystem", 36, WHITE, True)

# Columns
add_shape(slide, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5), BG_CARD)
add_text(slide, Inches(1.1), Inches(2.1), Inches(3.1), Inches(0.5), "Frontend", 22, ACCENT, True)
items = ["React 19 (Vite)", "Lucide Icons", "Recharts Viz", "Vanilla CSS", "Glassmorphism"]
add_bullet_text(slide, Inches(1.1), Inches(2.8), Inches(3.1), Inches(3.5), items, 15, LIGHT_GRAY, "›")

add_shape(slide, Inches(4.85), Inches(1.8), Inches(3.7), Inches(5), BG_CARD)
add_text(slide, Inches(5.15), Inches(2.1), Inches(3.1), Inches(0.5), "Backend", 22, SUCCESS, True)
items = ["Node.js / Express", "MongoDB Atlas", "Nodemailer", "WebSpeech API", "RESTful Routes"]
add_bullet_text(slide, Inches(5.15), Inches(2.8), Inches(3.1), Inches(3.5), items, 15, LIGHT_GRAY, "›")

add_shape(slide, Inches(8.9), Inches(1.8), Inches(3.7), Inches(5), BG_CARD)
add_text(slide, Inches(9.2), Inches(2.1), Inches(3.1), Inches(0.5), "AI Services", 22, ACCENT2, True)
items = ["Google Gemini AI", "OpenAI Wrapper", "NLP Intent Parsing", "Emotion Analysis", "TTS Engines"]
add_bullet_text(slide, Inches(9.2), Inches(2.8), Inches(3.1), Inches(3.5), items, 15, LIGHT_GRAY, "›")

# ─── SLIDE 13: SYSTEM ARCHITECTURE ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
set_slide_transition(slide, "push")
add_accent_bar(slide, Inches(0.8), Inches(0.6), Inches(0.15), Inches(0.6), ACCENT)
add_text(slide, Inches(1.2), Inches(0.55), Inches(10), Inches(0.7), "System Architecture Flow", 36, WHITE, True)

# Visual flow
add_shape(slide, Inches(2), Inches(1.8), Inches(9.3), Inches(1), BG_CARD)
add_text(slide, Inches(2), Inches(2.0), Inches(9.3), Inches(0.6), "Voice / UI Input Layer", 22, WHITE, True, PP_ALIGN.CENTER)

add_text(slide, Inches(6.2), Inches(2.8), Inches(1), Inches(0.5), "▼", 30, ACCENT, False, PP_ALIGN.CENTER)

add_shape(slide, Inches(2), Inches(3.3), Inches(9.3), Inches(1.5), BG_CARD)
add_accent_bar(slide, Inches(2), Inches(3.3), Inches(9.3), Pt(4), ACCENT2)
add_text(slide, Inches(2.5), Inches(3.6), Inches(8.3), Inches(0.5), "Core Engine (React + Node.js)", 20, ACCENT2, True, PP_ALIGN.CENTER)
add_text(slide, Inches(2.5), Inches(4.1), Inches(8.3), Inches(0.5), "Intent Parsing  •  AI Logic  •  Data Processing  •  Email Services", 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)

add_text(slide, Inches(6.2), Inches(4.8), Inches(1), Inches(0.5), "▼", 30, SUCCESS, False, PP_ALIGN.CENTER)

add_shape(slide, Inches(3.5), Inches(5.5), Inches(6.3), Inches(1), BG_CARD)
add_text(slide, Inches(3.5), Inches(5.7), Inches(6.3), Inches(0.6), "MongoDB Atlas & AI Cloud", 22, SUCCESS, True, PP_ALIGN.CENTER)

# ─── SLIDE 14: COMPETITIVE ADVANTAGE ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
set_slide_transition(slide, "push")
add_accent_bar(slide, Inches(0.8), Inches(0.6), Inches(0.15), Inches(0.6), ACCENT)
add_text(slide, Inches(1.2), Inches(0.55), Inches(10), Inches(0.7), "Competitive Advantage", 36, WHITE, True)

diffs = [
    ("⚡ Speed", "Voice-first navigation reduces CRM admin time by over 50%."),
    ("🧠 Intelligence", "Real-time AI insights that predict win rates and highlight risks."),
    ("🔗 Unity", "All-in-one platform eliminating the need for 5 separate tools."),
    ("💎 Design", "Premium visuals that drive user adoption and elite performance.")
]
for i, (title, desc) in enumerate(diffs):
    y = Inches(1.8 + i * 1.3)
    add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(1.1), BG_CARD)
    add_accent_bar(slide, Inches(0.8), y, Inches(0.12), Inches(1.1), ACCENT)
    add_text(slide, Inches(1.3), y + Inches(0.2), Inches(4), Inches(0.5), title, 22, WHITE, True)
    add_text(slide, Inches(5), y + Inches(0.3), Inches(7), Inches(0.5), desc, 18, LIGHT_GRAY)

# ─── SLIDE 15: THANK YOU ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
set_slide_transition(slide, "fade")
add_circle(slide, Inches(-1.5), Inches(-1.5), Inches(5), ACCENT)
add_circle(slide, Inches(10), Inches(4), Inches(5), ACCENT2)
add_shape(slide, Inches(2.5), Inches(1.5), Inches(8.3), Inches(4.5), BG_CARD)
add_accent_bar(slide, Inches(5.5), Inches(2.2), Inches(2.3), Pt(5), ACCENT)
add_text(slide, Inches(3), Inches(2.5), Inches(7.3), Inches(1), "Thank You", 48, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(3), Inches(3.5), Inches(7.3), Inches(0.6), "AuraCRM — The Future of Sales", 22, ACCENT2, False, PP_ALIGN.CENTER)
add_text(slide, Inches(3), Inches(4.6), Inches(7.3), Inches(0.6), "Ready to Revolutionize Your Pipeline?", 18, LIGHT_GRAY, False, PP_ALIGN.CENTER)

output = os.path.join(os.path.dirname(__file__), "AuraCRM_Full_Presentation.pptx")
prs.save(output)
print(f"✅ Full Presentation with Animations saved to: {output}")
